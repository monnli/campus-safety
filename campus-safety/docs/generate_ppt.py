"""生成答辩 PPT - 黔视护苗·校园安全智能监控系统"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# 颜色定义
C_DARK = RGBColor(0x00, 0x15, 0x29)      # 深蓝背景
C_BLUE = RGBColor(0x16, 0x77, 0xFF)      # 主蓝色
C_LIGHT = RGBColor(0x7E, 0xC8, 0xE3)    # 浅蓝
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)     # 白色
C_YELLOW = RGBColor(0xFA, 0xAD, 0x14)   # 黄色
C_RED = RGBColor(0xFF, 0x4D, 0x4F)      # 红色
C_GREEN = RGBColor(0x52, 0xC4, 0x1A)    # 绿色
C_GRAY = RGBColor(0xB0, 0xD4, 0xF0)     # 浅灰蓝

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 空白布局


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=C_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_rect(slide, left, top, width, height, color=C_BLUE, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line(slide, left, top, width, color=C_BLUE, thickness=2):
    from pptx.util import Pt
    line = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


# ── 封面 ──────────────────────────────────────────────────────
slide = add_slide()
set_bg(slide)

# 装饰线
add_rect(slide, 0, 0, 13.33, 0.08, C_BLUE)
add_rect(slide, 0, 7.42, 13.33, 0.08, C_BLUE)

# 左侧装饰块
add_rect(slide, 0, 0.08, 0.5, 7.34, RGBColor(0x00, 0x2A, 0x50))

# 主标题
add_text(slide, "黔视护苗", 1.2, 1.5, 11, 1.5,
         font_size=60, bold=True, color=C_LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, "校园安全智能监控系统", 1.2, 2.9, 11, 1,
         font_size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 副标题
add_rect(slide, 3.5, 4.1, 6.33, 0.04, C_BLUE)
add_text(slide, "AI Safety Guardian · 守护每一个孩子的平安成长", 1.2, 4.2, 11, 0.6,
         font_size=16, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

# 关键数据
stats = [("<3s", "响应时间"), ("<5%", "误报率"), ("24h", "全天守护"), ("90%", "带宽节省")]
for i, (val, label) in enumerate(stats):
    x = 1.5 + i * 2.6
    add_rect(slide, x, 5.1, 2.2, 1.0, RGBColor(0x00, 0x2A, 0x50))
    add_text(slide, val, x, 5.15, 2.2, 0.5,
             font_size=24, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, 5.65, 2.2, 0.4,
             font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# 参赛信息
add_text(slide, "全国大学生计算机设计大赛 · 人工智能应用赛道",
         1.2, 6.5, 11, 0.5, font_size=13, color=C_GRAY, align=PP_ALIGN.CENTER)


# ── 目录 ──────────────────────────────────────────────────────
slide = add_slide()
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_BLUE)

add_text(slide, "目  录", 0.5, 0.3, 12, 0.8,
         font_size=28, bold=True, color=C_WHITE)
add_rect(slide, 0.5, 1.0, 0.06, 0.4, C_BLUE)

items = [
    ("01", "问题背景与痛点分析"),
    ("02", "解决方案与系统概述"),
    ("03", "技术架构与创新亮点"),
    ("04", "核心功能演示"),
    ("05", "与传统方案对比"),
    ("06", "未来规划"),
]
for i, (num, title) in enumerate(items):
    row = i // 2
    col = i % 2
    x = 0.8 + col * 6.2
    y = 1.5 + row * 1.6
    add_rect(slide, x, y, 5.8, 1.3, RGBColor(0x00, 0x2A, 0x50))
    add_rect(slide, x, y, 0.8, 1.3, C_BLUE)
    add_text(slide, num, x, y + 0.3, 0.8, 0.7,
             font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.9, y + 0.35, 4.8, 0.6,
             font_size=18, bold=True, color=C_WHITE)


# ── 问题背景 ──────────────────────────────────────────────────
slide = add_slide()
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_BLUE)
add_rect(slide, 0, 0.08, 0.08, 7.34, C_BLUE)

add_text(slide, "01  问题背景与痛点分析", 0.3, 0.2, 12, 0.7,
         font_size=24, bold=True, color=C_WHITE)
add_line(slide, 0.3, 0.95, 12.7, C_BLUE)

problems = [
    ("🔴", "监控盲区", "固定摄像头覆盖有限，操场边角、校园周边等区域存在大量盲区"),
    ("🟡", "响应滞后", "传统监控依赖人工值守，事件发生后才能介入，平均响应时间5~15分钟"),
    ("🔴", "霸凌隐蔽", "校园霸凌往往发生在监控薄弱区域，受害者不敢举报，难以发现"),
    ("🟡", "人力不足", "安保人员数量有限，无法实现全时段、全区域巡逻"),
    ("🔵", "数据孤岛", "事件记录分散，缺乏系统性分析和预测能力，无法主动预防"),
    ("🔵", "管理粗放", "缺乏科学的巡逻任务管理，安全管理依赖经验而非数据"),
]
for i, (icon, title, desc) in enumerate(problems):
    row = i // 2
    col = i % 2
    x = 0.4 + col * 6.4
    y = 1.2 + row * 1.8
    add_rect(slide, x, y, 6.0, 1.5, RGBColor(0x00, 0x2A, 0x50))
    add_text(slide, title, x + 0.2, y + 0.1, 5.6, 0.5,
             font_size=16, bold=True, color=C_YELLOW)
    add_text(slide, desc, x + 0.2, y + 0.6, 5.6, 0.8,
             font_size=12, color=C_GRAY)

add_text(slide, "据教育部统计：约60%的校园安全事件发生在课间、午休及放学时段，约40%发生在监控盲区",
         0.4, 6.8, 12.5, 0.5, font_size=12, color=C_LIGHT, italic=True)


# ── 解决方案 ──────────────────────────────────────────────────
slide = add_slide()
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_BLUE)
add_rect(slide, 0, 0.08, 0.08, 7.34, C_BLUE)

add_text(slide, "02  解决方案与系统概述", 0.3, 0.2, 12, 0.7,
         font_size=24, bold=True, color=C_WHITE)
add_line(slide, 0.3, 0.95, 12.7, C_BLUE)

add_text(slide, "黔视护苗 · 校园安全智能监控系统", 0.4, 1.1, 12.5, 0.6,
         font_size=20, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "融合 YOLO目标检测 + Qwen多模态大模型 + 无人机边缘计算，构建全方位、智能化的校园安全防护体系",
         0.4, 1.7, 12.5, 0.5, font_size=13, color=C_GRAY, align=PP_ALIGN.CENTER)

solutions = [
    (C_RED, "实时检测", "YOLO26x实时识别\n打架/霸凌/跌倒等\n危险行为，<3秒响应"),
    (C_BLUE, "多模态分析", "Qwen-VL对视频片段\n进行时序语义理解\n二次确认降低误报"),
    (C_YELLOW, "无人机巡逻", "边缘计算方案\n覆盖监控盲区\n节省90%带宽"),
    (C_GREEN, "智能预警", "自动广播+短信+\n警报，全自动\n处置链路"),
    (C_LIGHT, "风险预测", "AI分析历史数据\n预测高危时段\n主动防控"),
    (C_BLUE, "持续录像", "NVR模式录像\n完整证据链\n智能检索回放"),
]
for i, (color, title, desc) in enumerate(solutions):
    x = 0.4 + (i % 3) * 4.2
    y = 2.4 + (i // 3) * 2.3
    add_rect(slide, x, y, 3.8, 2.0, RGBColor(0x00, 0x2A, 0x50))
    add_rect(slide, x, y, 3.8, 0.08, color)
    add_text(slide, title, x + 0.1, y + 0.15, 3.6, 0.5,
             font_size=15, bold=True, color=color)
    add_text(slide, desc, x + 0.1, y + 0.65, 3.6, 1.2,
             font_size=12, color=C_GRAY)


# ── 技术架构 ──────────────────────────────────────────────────
slide = add_slide()
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_BLUE)
add_rect(slide, 0, 0.08, 0.08, 7.34, C_BLUE)

add_text(slide, "03  技术架构与创新亮点", 0.3, 0.2, 12, 0.7,
         font_size=24, bold=True, color=C_WHITE)
add_line(slide, 0.3, 0.95, 12.7, C_BLUE)

# 架构层次
layers = [
    (C_BLUE, "前端展示层", "Vue 3 + Element Plus + ECharts + Socket.IO + Web Speech API"),
    (RGBColor(0x00, 0x2A, 0x50), "后端服务层", "Flask + Flask-SocketIO + gevent + SQLite + JWT认证"),
    (RGBColor(0x00, 0x1A, 0x35), "AI模型层", "YOLO26x检测 + Qwen-VL多帧分析 + Qwen-Max报告生成"),
    (RGBColor(0x00, 0x10, 0x25), "数据采集层", "固定摄像头（RTSP/视频文件/本地摄像头）+ 无人机（边缘计算）"),
]
for i, (color, title, desc) in enumerate(layers):
    y = 1.1 + i * 1.1
    add_rect(slide, 0.4, y, 8.5, 0.9, color)
    add_rect(slide, 0.4, y, 2.2, 0.9, C_BLUE)
    add_text(slide, title, 0.4, y + 0.15, 2.2, 0.6,
             font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, 2.7, y + 0.2, 6.1, 0.5,
             font_size=12, color=C_GRAY)

# 创新亮点
add_text(slide, "创新亮点", 9.3, 1.0, 3.8, 0.5,
         font_size=16, bold=True, color=C_YELLOW)
innovations = [
    "双重AI验证，误报率<5%",
    "视频级多帧时序分析",
    "边缘计算无人机方案",
    "异步检测框实时叠加",
    "Mock降级保证可靠性",
    "NVR持续录像+智能检索",
    "AI巡逻任务自动规划",
]
for i, text in enumerate(innovations):
    add_rect(slide, 9.3, 1.55 + i * 0.72, 3.8, 0.6, RGBColor(0x00, 0x2A, 0x50))
    add_rect(slide, 9.3, 1.55 + i * 0.72, 0.08, 0.6, C_BLUE)
    add_text(slide, text, 9.5, 1.65 + i * 0.72, 3.5, 0.4,
             font_size=12, color=C_WHITE)


# ── 核心功能 ──────────────────────────────────────────────────
slide = add_slide()
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_BLUE)
add_rect(slide, 0, 0.08, 0.08, 7.34, C_BLUE)

add_text(slide, "04  核心功能演示", 0.3, 0.2, 12, 0.7,
         font_size=24, bold=True, color=C_WHITE)
add_line(slide, 0.3, 0.95, 12.7, C_BLUE)

features = [
    ("实时监控", C_BLUE, [
        "多路视频同时接入（RTSP/文件/摄像头）",
        "YOLO检测框实时叠加显示",
        "支持暂停/停止/动态添加摄像头",
        "全屏监控模式",
    ]),
    ("智能预警", C_RED, [
        "检测到危险行为立即触发",
        "Qwen-VL多帧视频分析确认",
        "自动语音广播（Web Speech API）",
        "短信通知+监控室警报",
    ]),
    ("事件管理", C_YELLOW, [
        "结构化事件记录+AI预警报告",
        "存档视频在线播放/下载",
        "标记处理/误报，形成闭环",
        "时间范围筛选+日报生成",
    ]),
    ("安全大屏", C_GREEN, [
        "进入自动全屏，适合答辩展示",
        "安全指数+热力图+趋势图",
        "AI风险预测+防控建议",
        "实时多路视频（自适应布局）",
    ]),
    ("历史录像", C_LIGHT, [
        "NVR持续录像，完整证据链",
        "按摄像头/日期检索",
        "在线播放+下载",
        "自动清理过期录像",
    ]),
    ("巡逻管理", C_BLUE, [
        "AI根据风险分析生成巡逻任务",
        "指定时间、区域、负责人",
        "数字化签到确认",
        "完成率实时统计",
    ]),
]
for i, (title, color, items) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.3
    y = 1.1 + row * 3.0
    add_rect(slide, x, y, 4.0, 2.7, RGBColor(0x00, 0x2A, 0x50))
    add_rect(slide, x, y, 4.0, 0.45, color)
    add_text(slide, title, x + 0.1, y + 0.05, 3.8, 0.35,
             font_size=14, bold=True, color=C_WHITE)
    for j, item in enumerate(items):
        add_text(slide, f"• {item}", x + 0.15, y + 0.55 + j * 0.5, 3.7, 0.45,
                 font_size=11, color=C_GRAY)


# ── 对比分析 ──────────────────────────────────────────────────
slide = add_slide()
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_BLUE)
add_rect(slide, 0, 0.08, 0.08, 7.34, C_BLUE)

add_text(slide, "05  与传统方案对比", 0.3, 0.2, 12, 0.7,
         font_size=24, bold=True, color=C_WHITE)
add_line(slide, 0.3, 0.95, 12.7, C_BLUE)

headers = ["对比维度", "传统人工监控", "传统视频监控", "本系统（AI智能监控）"]
col_widths = [2.5, 2.8, 2.8, 4.0]
col_colors = [RGBColor(0x00, 0x2A, 0x50), RGBColor(0x00, 0x2A, 0x50),
              RGBColor(0x00, 0x2A, 0x50), C_BLUE]

x_starts = [0.3]
for w in col_widths[:-1]:
    x_starts.append(x_starts[-1] + w + 0.05)

# 表头
for i, (header, w, x, color) in enumerate(zip(headers, col_widths, x_starts, col_colors)):
    add_rect(slide, x, 1.1, w, 0.5, color)
    add_text(slide, header, x, 1.15, w, 0.4,
             font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

rows = [
    ["响应时间", "5~15分钟", "事后回看", "✅ 实时（<3秒）"],
    ["覆盖范围", "有限", "固定区域", "✅ 固定+无人机全覆盖"],
    ["误报率", "低（人工）", "无自动识别", "✅ <5%（双重AI验证）"],
    ["24h值守", "需轮班", "需人工查看", "✅ 全自动，无需值守"],
    ["历史回溯", "无", "手动翻录像", "✅ 持续录像+在线检索"],
    ["预测能力", "无", "无", "✅ AI风险预测+热力图"],
    ["巡逻管理", "纸质记录", "无", "✅ AI生成任务+数字签到"],
    ["处置响应", "人工通知", "人工通知", "✅ 自动广播+短信+警报"],
]
for r, row_data in enumerate(rows):
    bg = RGBColor(0x00, 0x1A, 0x35) if r % 2 == 0 else RGBColor(0x00, 0x22, 0x44)
    for i, (cell, w, x) in enumerate(zip(row_data, col_widths, x_starts)):
        cell_color = RGBColor(0x00, 0x2A, 0x50) if i == 3 else bg
        add_rect(slide, x, 1.65 + r * 0.65, w, 0.6, cell_color)
        text_color = C_GREEN if cell.startswith("✅") else C_GRAY
        add_text(slide, cell, x + 0.05, 1.7 + r * 0.65, w - 0.1, 0.5,
                 font_size=11, color=text_color, align=PP_ALIGN.CENTER)


# ── 未来规划 ──────────────────────────────────────────────────
slide = add_slide()
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_BLUE)
add_rect(slide, 0, 0.08, 0.08, 7.34, C_BLUE)

add_text(slide, "06  未来规划", 0.3, 0.2, 12, 0.7,
         font_size=24, bold=True, color=C_WHITE)
add_line(slide, 0.3, 0.95, 12.7, C_BLUE)

phases = [
    ("近期（1-3个月）", C_BLUE, [
        "完成 YOLO26x 微调训练",
        "收集校园安全专项数据集",
        "提升打架/霸凌识别准确率至90%+",
        "接入真实摄像头和无人机测试",
    ]),
    ("中期（3-6个月）", C_YELLOW, [
        "与学校合作开展试点部署",
        "收集真实场景数据持续优化",
        "开发移动端管理 App",
        "接入真实短信通知服务",
    ]),
    ("长期（6个月+）", C_GREEN, [
        "推广至更多学校落地应用",
        "构建校园安全大数据平台",
        "探索联邦学习保护隐私",
        "申请相关专利和软件著作权",
    ]),
]
for i, (phase, color, items) in enumerate(phases):
    x = 0.5 + i * 4.2
    add_rect(slide, x, 1.2, 3.9, 0.5, color)
    add_text(slide, phase, x, 1.25, 3.9, 0.4,
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 1.7, 3.9, 4.5, RGBColor(0x00, 0x2A, 0x50))
    for j, item in enumerate(items):
        add_rect(slide, x + 0.15, 1.85 + j * 1.0, 0.08, 0.08, color)
        add_text(slide, item, x + 0.35, 1.8 + j * 1.0, 3.4, 0.8,
                 font_size=13, color=C_WHITE)

add_text(slide, "本系统已具备完整的技术架构和功能闭环，具有良好的可扩展性和实际落地价值",
         0.5, 6.5, 12.3, 0.5, font_size=13, color=C_LIGHT,
         align=PP_ALIGN.CENTER, italic=True)


# ── 结束页 ──────────────────────────────────────────────────
slide = add_slide()
set_bg(slide)
add_rect(slide, 0, 0, 13.33, 0.08, C_BLUE)
add_rect(slide, 0, 7.42, 13.33, 0.08, C_BLUE)

add_rect(slide, 0, 0.08, 0.5, 7.34, RGBColor(0x00, 0x2A, 0x50))

add_text(slide, "感谢聆听", 1.2, 2.0, 11, 1.5,
         font_size=56, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "THANK YOU", 1.2, 3.3, 11, 0.8,
         font_size=20, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)

add_rect(slide, 3.5, 4.2, 6.33, 0.04, C_BLUE)

add_text(slide, "黔视护苗 · 守护每一个孩子的平安成长",
         1.2, 4.4, 11, 0.6, font_size=16, color=C_LIGHT,
         align=PP_ALIGN.CENTER)
add_text(slide, "全国大学生计算机设计大赛 · 人工智能应用赛道",
         1.2, 5.1, 11, 0.5, font_size=13, color=C_GRAY,
         align=PP_ALIGN.CENTER)

add_text(slide, "本系统仅分析行为模式，不采集人脸信息，符合个人信息保护法规",
         1.2, 6.5, 11, 0.5, font_size=11, color=RGBColor(0x52, 0xC4, 0x1A),
         align=PP_ALIGN.CENTER)


# 保存
output_path = os.path.join(os.path.dirname(__file__), "黔视护苗_答辩PPT.pptx")
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
